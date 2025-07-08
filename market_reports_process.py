import os
import json
import re
import traceback
import logging
import requests
from flask import Flask, request, jsonify
from docxtpl import DocxTemplate
from pptx import Presentation
from pptx.util import Inches
from drive_utils import upload_to_drive

# Templates directory
TEMPLATES_DIR = os.path.join(os.path.dirname(__file__), "templates")
DOCX_TEMPLATE = os.path.join(TEMPLATES_DIR, "Market_Gap_Analysis_Template.docx")
PPTX_TEMPLATE = os.path.join(TEMPLATES_DIR, "Market_Gap_Analysis_Template.pptx")

# Configure logging
logging.basicConfig(level=logging.INFO)

app = Flask(__name__)


def to_direct_drive_url(url: str) -> str:
    """
    Convert a Google Drive share URL into a direct-download URL.
    """
    m = re.search(r"[?&]id=([^&]+)", url)
    if m:
        return f"https://drive.google.com/uc?export=download&id={m.group(1)}"
    m = re.search(r"/d/([^/]+)", url)
    if m:
        return f"https://drive.google.com/uc?export=download&id={m.group(1)}"
    return url


def download_chart(url: str, local_path: str) -> str:
    """
    Download a chart PNG from Drive and save locally.
    """
    direct_url = to_direct_drive_url(url)
    os.makedirs(os.path.dirname(local_path), exist_ok=True)
    resp = requests.get(direct_url)
    resp.raise_for_status()
    with open(local_path, "wb") as f:
        f.write(resp.content)
    return local_path


def replace_placeholder(slide, key: str, text: str):
    """
    Replace all occurrences of {{key}} in slide text frames with the given text.
    """
    pattern = re.compile(rf"\{{{{\s*{key}\s*}}}}")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Combine runs into a single text blob
        full_text = ''.join(run.text for run in shape.text_frame.paragraphs[0].runs)
        if pattern.search(full_text):
            new_text = pattern.sub(text, full_text)
            # Clear and set as single run
            frame = shape.text_frame
            frame.clear()
            p = frame.add_paragraph()
            p.text = new_text


@app.route('/start_market_gap', methods=['POST'])
def start_market_gap():
    data = request.get_json(force=True)
    logging.info("📦 Incoming payload:\n%s", json.dumps(data, indent=2))

    session_id = data.get("session_id")
    email = data.get("email", "")
    folder_id = data.get("folder_id", "")
    if not session_id:
        return jsonify({"error": "Missing session_id"}), 400

    local_path = os.path.join("temp_sessions", session_id)
    os.makedirs(local_path, exist_ok=True)

    try:
        result = generate_market_reports(
            session_id=session_id,
            email=email,
            folder_id=folder_id,
            payload=data,
            local_path=local_path
        )
        
def generate_market_reports(session_id: str,
                            email: str,
                            folder_id: str,
                            payload: dict,
                            local_path: str) -> dict:
    """
    Renders DOCX and PPTX using templates, uploads to Drive, and constructs the result payload.
    """
    # 1. Generate DOCX report
    doc = DocxTemplate(DOCX_TEMPLATE)
    context = payload.get("content", {}).copy()
    context["date"] = payload.get("date", "")
    context["organization_name"] = payload.get("organization_name", "")

    docx_filename = f"market_gap_analysis_report_{session_id}.docx"
    docx_path = os.path.join(local_path, docx_filename)
    doc.render(context)
    doc.save(docx_path)
    docx_url = upload_to_drive(docx_path, session_id, folder_id)

    # 2. Generate PPTX executive report
    pres = Presentation(PPTX_TEMPLATE)
    content = payload.get("content", {})
    charts = payload.get("charts", {})

    # Replace text placeholders on specific slides
    content_slide_map = {
        1: "executive_summary",
        2: "current_state_overview",
        3: "hardware_gap_analysis",
        4: "software_gap_analysis",
        5: "market_benchmarking",
    }
    for idx, key in content_slide_map.items():
        if idx < len(pres.slides):
            replace_placeholder(pres.slides[idx], key, content.get(key, ""))

    # Insert charts into named picture placeholders
    chart_placeholders = {
        "hardware_tier": charts.get("hardware_insights_tier"),
        "software_tier": charts.get("software_insights_tier"),
    }
    for slide in pres.slides:
        for ph_name, chart_url in chart_placeholders.items():
            if chart_url:
                local_chart = download_chart(chart_url, os.path.join(local_path, f"{ph_name}.png"))
                if os.path.exists(local_chart):
                    for ph in slide.placeholders:
                        if ph.name == ph_name:
                            ph.insert_picture(local_chart)
                            break

    pptx_filename = f"market_gap_analysis_executive_report_{session_id}.pptx"
    pptx_path = os.path.join(local_path, pptx_filename)
    pres.save(pptx_path)
    pptx_url = upload_to_drive(pptx_path, session_id, folder_id)

    # 3. Construct result payload
    result = {
        "session_id": session_id,
        "gpt_module": "gap_market",
        "status": "complete",
        "content": content,
        "charts": charts,
        "files": [{"file_name": f["file_name"], "file_url": f["file_url"]} for f in payload.get("files", [])],
        "file_1_name": os.path.basename(docx_path),
        "file_1_url": docx_url,
        "file_2_name": os.path.basename(pptx_path),
        "file_2_url": pptx_url
    }

    return result

    # Invoke IT Summarizer endpoint
        summarizer_url = "https://it-summarizer-api.onrender.com"
        try:
            requests.post(summarizer_url, json=result, timeout=30)
        except Exception:
            pass

        return jsonify(result), 200
    except Exception as e:
        logging.error(f"🔥 Market Reports generation failed: {e}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    port = int(os.getenv("PORT", 10000))
    app.run(host="0.0.0.0", port=port)
